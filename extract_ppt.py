import pptx
import os

def extract_pptx_text(pptx_path, output_path):
    prs = pptx.Presentation(pptx_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        text_runs.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(text_runs))

if __name__ == "__main__":
    pptx_file = "第1章 绪论、骨学.pptx"
    output_file = "extracted_content.txt"
    if os.path.exists(pptx_file):
        extract_pptx_text(pptx_file, output_file)
        print(f"Extraction successful. Saved to {output_file}")
    else:
        print(f"File {pptx_file} not found.")
